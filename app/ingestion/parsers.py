"""🛠️ Smart Parsing (LOCS/02_INGESTION_ENGINE.md §1).

All parsing runs entirely on-device — no external OCR service or cloud API.
Every function takes a file path and returns the extracted plain text.
"""
from __future__ import annotations

import logging
from pathlib import Path

log = logging.getLogger(__name__)


def parse_pdf(path: Path) -> str:
    """Parse a PDF locally.

    Primary: pypdf (fast, good for text-based PDFs).
    Fallback: pdfplumber (slower but better at complex layouts/tables).
    Handles multi-page PDFs with no page-count limits.
    """
    from pypdf import PdfReader

    try:
        reader = PdfReader(str(path))
        pages = [page.extract_text() or "" for page in reader.pages]
        text = "\n\n".join(p for p in pages if p.strip())
        if text.strip():
            return text
        log.warning("pypdf returned empty text for %s — falling back to pdfplumber", path.name)
    except Exception as exc:  # noqa: BLE001 — parser fallback should never kill ingestion
        log.warning("pypdf failed on %s (%s) — falling back to pdfplumber", path.name, exc)

    # Fallback for complex layouts
    import pdfplumber

    with pdfplumber.open(str(path)) as pdf:
        pages = []
        for page in pdf.pages:
            extracted = page.extract_text() or ""
            tables = page.extract_tables()
            for table in tables:
                for row in table:
                    extracted += "\n" + " | ".join(cell or "" for cell in row)
            pages.append(extracted)
        return "\n\n".join(p for p in pages if p.strip())


def parse_html(path: Path) -> str:
    """Parse HTML via BeautifulSoup — strips script/style/meta, keeps readable text."""
    from bs4 import BeautifulSoup

    with path.open("r", encoding="utf-8", errors="ignore") as fh:
        soup = BeautifulSoup(fh.read(), "html.parser")

    for tag in soup(["script", "style", "meta", "noscript", "header", "footer"]):
        tag.decompose()

    text = soup.get_text(separator="\n", strip=True)
    return text


def parse_docx(path: Path) -> str:
    """Parse .docx via python-docx."""
    import docx

    document = docx.Document(str(path))
    parts = [p.text for p in document.paragraphs if p.text.strip()]
    for table in document.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text.strip() for cell in row.cells))
    return "\n\n".join(parts)


def parse_pptx(path: Path) -> str:
    """Parse .pptx via python-pptx (text from every slide)."""
    from pptx import Presentation

    prs = Presentation(str(path))
    parts = []
    for idx, slide in enumerate(prs.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        if slide_text:
            parts.append(f"--- Slide {idx} ---\n" + "\n".join(slide_text))
    return "\n\n".join(parts)


def parse_text(path: Path) -> str:
    """Plain text loader."""
    return path.read_text(encoding="utf-8", errors="ignore")


# Extension → parser mapping. "Universal" ingestion just looks this up.
PARSERS = {
    ".pdf": parse_pdf,
    ".html": parse_html,
    ".htm": parse_html,
    ".txt": parse_text,
    ".md": parse_text,
    ".docx": parse_docx,
    ".pptx": parse_pptx,
}


def parse_file(path: Path) -> str:
    """Dispatch a file to the right parser. Returns "" for unsupported types."""
    parser = PARSERS.get(path.suffix.lower())
    if parser is None:
        log.warning("Unsupported file type: %s", path)
        return ""
    return parser(path)
